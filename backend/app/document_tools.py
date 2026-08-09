from __future__ import annotations

import io
import json
from typing import Iterable

import fitz
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation

from .artifacts import save_artifact


def docx_to_text(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    blocks = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            blocks.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(blocks)


def xlsx_to_text(data: bytes, max_rows_per_sheet: int = 5000) -> str:
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows_per_sheet:
                parts.append("[row limit reached]")
                break
            parts.append(" | ".join("" if v is None else str(v) for v in row))
    return "\n".join(parts)


def pptx_to_text(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def create_pdf(title: str, text: str) -> dict:
    doc = fitz.open()
    page = doc.new_page()
    y = 60
    page.insert_text((50, y), title[:120], fontsize=18)
    y += 35
    for paragraph in text.splitlines() or [""]:
        if y > 760:
            page = doc.new_page(); y = 60
        rect = fitz.Rect(50, y, 545, min(y + 140, 790))
        spare = page.insert_textbox(rect, paragraph, fontsize=10, lineheight=1.25)
        used = max(18, 140 - max(spare, 0))
        y += used + 8
    data = doc.tobytes(garbage=4, deflate=True)
    return save_artifact(f"{title or 'document'}.pdf", data, "application/pdf", {"tool": "create_pdf"})


def merge_pdfs(files: Iterable[bytes], name: str = "merged.pdf") -> dict:
    out = fitz.open()
    count = 0
    for raw in files:
        src = fitz.open(stream=raw, filetype="pdf")
        out.insert_pdf(src)
        count += src.page_count
    data = out.tobytes(garbage=4, deflate=True)
    return save_artifact(name, data, "application/pdf", {"tool": "merge_pdf", "pages": count})


def annotate_pdf(data: bytes, page_number: int, text: str, x: float = 50, y: float = 50) -> dict:
    doc = fitz.open(stream=data, filetype="pdf")
    if page_number < 1 or page_number > doc.page_count:
        raise ValueError("page_number is outside the PDF")
    page = doc[page_number - 1]
    annot = page.add_text_annot((x, y), text)
    annot.update()
    output = doc.tobytes(garbage=4, deflate=True)
    return save_artifact("annotated.pdf", output, "application/pdf", {"tool": "annotate_pdf", "page": page_number})


def create_docx(title: str, text: str) -> dict:
    doc = Document()
    if title:
        doc.add_heading(title, level=1)
    for paragraph in text.split("\n\n"):
        doc.add_paragraph(paragraph)
    buf = io.BytesIO(); doc.save(buf)
    return save_artifact(f"{title or 'document'}.docx", buf.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", {"tool": "create_docx"})


def create_xlsx(rows: list[dict], name: str = "data.xlsx") -> dict:
    wb = Workbook(); ws = wb.active; ws.title = "Data"
    keys = list(rows[0].keys()) if rows else []
    if keys:
        ws.append(keys)
        for row in rows:
            ws.append([row.get(k) for k in keys])
    buf = io.BytesIO(); wb.save(buf)
    return save_artifact(name, buf.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", {"tool": "create_xlsx", "rows": len(rows)})
